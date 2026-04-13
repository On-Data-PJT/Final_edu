from __future__ import annotations

import logging
import math
import re
import tempfile
import time
import uuid
from dataclasses import dataclass, field
from datetime import date, datetime
from pathlib import Path
from urllib.parse import parse_qs, urlparse

try:
    from openai import OpenAI
except ImportError:  # pragma: no cover - handled at runtime when dependency is missing.
    OpenAI = None
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader
from yt_dlp import YoutubeDL
from youtube_transcript_api import YouTubeTranscriptApi

from final_edu.config import Settings, get_settings
from final_edu.models import RawTextSegment, SourceAsset, UploadedAsset
from final_edu.storage import ObjectStorage, create_object_storage
from final_edu.utils import format_seconds, normalize_text
from final_edu.youtube_cache import (
    YOUTUBE_STALE_TRANSCRIPT_WARNING,
    YOUTUBE_STT_BUDGET_EXCEEDED_WARNING,
    YOUTUBE_STT_DISABLED_WARNING,
    YOUTUBE_STT_FALLBACK_WARNING,
    YOUTUBE_STT_FILE_TOO_LARGE_WARNING,
    YoutubeCache,
    build_youtube_scraperapi_http_client,
    is_youtube_request_limited_error,
    is_youtube_stt_fallback_error,
    mark_youtube_request_limited,
    summarize_youtube_fetch_error,
    throttle_youtube_requests,
)

YOUTUBE_ID_RE = re.compile(r"(?:v=|youtu\.be/|shorts/)([0-9A-Za-z_-]{11})")
EXCEL_VOC_HEADER_HINTS = (
    "comment",
    "comments",
    "feedback",
    "review",
    "opinion",
    "response",
    "suggestion",
    "rating",
    "score",
    "week",
    "lecture",
    "class",
    "강의",
    "평가",
    "후기",
    "의견",
    "소감",
    "피드백",
    "불만",
    "개선",
    "만족",
    "건의",
    "응답",
    "주차",
    "차시",
)
EXCEL_VOC_TEXT_HEADER_HINTS = (
    "comment",
    "comments",
    "feedback",
    "review",
    "opinion",
    "response",
    "suggestion",
    "강의",
    "평가",
    "후기",
    "의견",
    "소감",
    "피드백",
    "불만",
    "개선",
    "건의",
    "응답",
)
EXCEL_VOC_META_HEADER_HINTS = (
    "week",
    "date",
    "lecture",
    "class",
    "rating",
    "score",
    "주차",
    "차시",
    "날짜",
    "만족도",
    "점수",
)
NUMERICISH_RE = re.compile(r"^[\d\s.,%+\-/:]+$")
TEXT_SIGNAL_RE = re.compile(r"[A-Za-z가-힣]")
VOC_SURVEY_QUESTION_ID_RE = re.compile(r"\b([AB]Q\d+(?:-\d+)?)\b", re.IGNORECASE)
VOC_SURVEY_FREE_TEXT_HINTS = (
    "기타 의견",
    "기타의견",
    "의견",
    "건의",
    "제안",
    "서술형",
    "자유롭게",
    "comment",
    "comments",
    "opinion",
    "feedback",
    "suggestion",
)
logger = logging.getLogger(__name__)


@dataclass(slots=True)
class TabularSheet:
    name: str
    rows: list[list[str]]


@dataclass(slots=True)
class VocSheetCandidate:
    sheet: TabularSheet
    score: float
    header: list[str]
    data_rows: list[list[str]]
    text_columns: list[int]
    meta_columns: list[int]
    response_like_rows: int
    numeric_ratio: float
    header_hint_count: int


@dataclass(slots=True)
class VocSurveyQuestionColumn:
    column_index: int
    question_id: str
    group: str
    label: str
    scale_max: int = 5


@dataclass(slots=True)
class VocSurveyCandidate:
    sheet: TabularSheet
    score: float
    header_row_index: int
    header: list[str]
    data_rows: list[list[str]]
    question_columns: list[VocSurveyQuestionColumn]
    free_text_columns: list[int]
    response_row_count: int


@dataclass(slots=True)
class VocExtractionResult:
    source: SourceAsset
    segments: list[RawTextSegment]
    warnings: list[str] = field(default_factory=list)
    response_count: int = 0
    question_scores: list[dict] = field(default_factory=list)


def _build_file_source_asset(asset: UploadedAsset, instructor_name: str) -> SourceAsset:
    suffix = asset.path.suffix.lower()
    return SourceAsset(
        id=uuid.uuid4().hex[:12],
        instructor_name=instructor_name,
        asset_type=suffix.lstrip(".") or "file",
        label=asset.original_name,
        origin=asset.original_name,
    )


def extract_file_asset(
    asset: UploadedAsset,
    instructor_name: str,
) -> tuple[SourceAsset, list[RawTextSegment], list[str]]:
    suffix = asset.path.suffix.lower()
    source = _build_file_source_asset(asset, instructor_name)

    if suffix == ".pdf":
        segments, warnings = _extract_pdf(asset.path, source, instructor_name)
    elif suffix == ".pptx":
        segments, warnings = _extract_pptx(asset.path, source, instructor_name)
    elif suffix == ".csv":
        segments, warnings = _extract_csv_file(asset.path, source, instructor_name)
    elif suffix in {".xlsx", ".xls"}:
        segments, warnings = _extract_excel_file(asset.path, source, instructor_name)
    elif suffix in {".txt", ".md"}:
        segments, warnings = _extract_text_file(asset.path, source, instructor_name)
    else:
        segments = []
        warnings = [f"{asset.original_name}: 지원하지 않는 파일 형식입니다."]

    source.warnings.extend(warnings)
    return source, segments, warnings


def extract_voc_asset(
    asset: UploadedAsset,
    instructor_name: str,
) -> VocExtractionResult:
    suffix = asset.path.suffix.lower()
    source = _build_file_source_asset(asset, instructor_name)

    if suffix == ".csv":
        return _extract_voc_csv_file(asset.path, source, instructor_name)
    if suffix in {".xlsx", ".xls"}:
        return _extract_voc_excel_file(asset.path, source, instructor_name)

    source, segments, warnings = extract_file_asset(asset, instructor_name)
    return VocExtractionResult(
        source=source,
        segments=segments,
        warnings=warnings,
        response_count=_estimate_tabular_response_count(source.asset_type, segments),
        question_scores=[],
    )


def extract_youtube_asset(
    url: str,
    instructor_name: str,
    settings: Settings | None = None,
    storage: ObjectStorage | None = None,
    allow_stt_fallback: bool = True,
) -> tuple[SourceAsset, list[RawTextSegment], list[str]]:
    settings = settings or get_settings()
    video_id = _extract_video_id(url)
    active_storage = storage or create_object_storage(settings)
    video_title = _lookup_cached_youtube_title(url, video_id, settings, active_storage)
    source = SourceAsset(
        id=uuid.uuid4().hex[:12],
        instructor_name=instructor_name,
        asset_type="youtube",
        label=video_title or f"YouTube {video_id or 'unknown'}",
        origin=url,
    )

    if not video_id:
        warning = f"{url}: 올바른 YouTube URL을 인식하지 못했습니다."
        source.warnings.append(warning)
        return source, [], [warning]

    warnings: list[str] = []
    try:
        transcript, fetch_warnings = _fetch_youtube_transcript(video_id, settings, storage=active_storage)
        warnings.extend(fetch_warnings)
    except Exception as exc:  # noqa: BLE001
        if is_youtube_stt_fallback_error(exc):
            if not allow_stt_fallback:
                warning = (
                    f"{url}: 준비 단계 샘플 자막 확인에서는 공개 자막을 찾지 못했습니다. "
                    "본분석에서는 STT fallback 대상이 될 수 있습니다."
                )
                source.warnings.append(warning)
                return source, [], [warning]
            stt_segments, stt_warnings = _try_youtube_stt_fallback(
                url,
                video_id,
                source,
                instructor_name,
                settings,
                storage=active_storage,
            )
            if stt_segments:
                source.warnings.extend(stt_warnings)
                return source, stt_segments, stt_warnings
            if stt_warnings:
                source.warnings.extend(stt_warnings)
                return source, [], stt_warnings
        warning = summarize_youtube_fetch_error(url, exc)
        source.warnings.append(warning)
        return source, [], [warning]

    segments: list[RawTextSegment] = []
    for item in transcript:
        text = normalize_text(_transcript_text(item))
        if not text:
            continue
        locator = format_seconds(_transcript_start(item))
        segments.append(
            RawTextSegment(
                source_id=source.id,
                instructor_name=instructor_name,
                source_label=source.label,
                source_type=source.asset_type,
                locator=locator,
                text=text,
            )
        )

    if not segments:
        warnings.append(f"{url}: 자막은 조회됐지만 실제 분석 가능한 텍스트가 없습니다.")

    source.warnings.extend(warnings)

    return source, segments, warnings


def _extract_pdf(
    path: Path,
    source: SourceAsset,
    instructor_name: str,
) -> tuple[list[RawTextSegment], list[str]]:
    warnings: list[str] = []
    reader = PdfReader(str(path))
    segments: list[RawTextSegment] = []
    empty_pages = 0

    for index, page in enumerate(reader.pages, start=1):
        text = normalize_text(page.extract_text() or "")
        if not text:
            empty_pages += 1
            continue
        segments.append(
            RawTextSegment(
                source_id=source.id,
                instructor_name=instructor_name,
                source_label=source.label,
                source_type=source.asset_type,
                locator=f"p.{index}",
                text=text,
            )
        )

    if empty_pages:
        warnings.append(
            f"{source.label}: 텍스트가 추출되지 않은 페이지 {empty_pages}개를 제외했습니다. "
            "스캔 PDF인 경우 정확도가 떨어질 수 있습니다."
        )
    if not segments:
        warnings.append(
            f"{source.label}: 분석 가능한 텍스트를 찾지 못했습니다. 텍스트형 PDF만 안정적으로 지원합니다."
        )

    return segments, warnings


def _extract_pptx(
    path: Path,
    source: SourceAsset,
    instructor_name: str,
) -> tuple[list[RawTextSegment], list[str]]:
    warnings: list[str] = []
    presentation = Presentation(str(path))
    segments: list[RawTextSegment] = []
    empty_slides = 0

    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            parts.extend(_shape_texts(shape))

        text = normalize_text(" ".join(part for part in parts if part))
        if not text:
            empty_slides += 1
            continue

        segments.append(
            RawTextSegment(
                source_id=source.id,
                instructor_name=instructor_name,
                source_label=source.label,
                source_type=source.asset_type,
                locator=f"s.{index}",
                text=text,
            )
        )

    if empty_slides:
        warnings.append(
            f"{source.label}: 텍스트가 없는 슬라이드 {empty_slides}개를 제외했습니다. "
            "이미지 중심 슬라이드는 반영되지 않을 수 있습니다."
        )
    if not segments:
        warnings.append(f"{source.label}: PPTX에서 분석 가능한 텍스트를 찾지 못했습니다.")

    return segments, warnings


def _extract_text_file(
    path: Path,
    source: SourceAsset,
    instructor_name: str,
) -> tuple[list[RawTextSegment], list[str]]:
    text = normalize_text(path.read_text(encoding="utf-8", errors="ignore"))
    if not text:
        warning = f"{source.label}: 텍스트 파일이 비어 있습니다."
        return [], [warning]

    segment = RawTextSegment(
        source_id=source.id,
        instructor_name=instructor_name,
        source_label=source.label,
        source_type="text",
        locator="full",
        text=text,
    )
    return [segment], []


def _extract_csv_file(
    path: Path,
    source: SourceAsset,
    instructor_name: str,
) -> tuple[list[RawTextSegment], list[str]]:
    import csv

    warnings: list[str] = []
    raw_text = path.read_text(encoding="utf-8", errors="ignore")
    rows = list(csv.reader(raw_text.splitlines()))
    if not rows:
        return [], [f"{source.label}: CSV 파일이 비어 있습니다."]

    header = rows[0] if rows else []
    data_rows = rows[1:] if len(rows) > 1 else []
    segments: list[RawTextSegment] = []

    for index, row in enumerate(data_rows or rows, start=1):
        row_text = normalize_text(" | ".join(cell for cell in row if cell))
        if not row_text:
            continue
        if header and data_rows:
            pairs = []
            for header_cell, value in zip(header, row):
                clean_header = normalize_text(header_cell)
                clean_value = normalize_text(value)
                if clean_value:
                    label = clean_header or "column"
                    pairs.append(f"{label}: {clean_value}")
            row_text = normalize_text(" | ".join(pairs)) or row_text
        segments.append(
            RawTextSegment(
                source_id=source.id,
                instructor_name=instructor_name,
                source_label=source.label,
                source_type="csv",
                locator=f"row.{index}",
                text=row_text,
            )
        )

    if not segments:
        warnings.append(f"{source.label}: CSV에서 분석 가능한 텍스트를 찾지 못했습니다.")

    return segments, warnings


def _extract_voc_csv_file(
    path: Path,
    source: SourceAsset,
    instructor_name: str,
) -> VocExtractionResult:
    import csv

    raw_text = path.read_text(encoding="utf-8", errors="ignore")
    rows = list(csv.reader(raw_text.splitlines()))
    if not rows:
        raise ValueError(f"{source.label}: CSV 파일이 비어 있습니다.")

    survey_candidate = _score_voc_survey_sheet_candidate(TabularSheet(name="csv", rows=rows))
    if survey_candidate is not None:
        return _build_voc_extraction_from_survey_candidate(
            source=source,
            instructor_name=instructor_name,
            survey_candidate=survey_candidate,
        )

    segments, warnings = _extract_csv_file(path, source, instructor_name)
    return VocExtractionResult(
        source=source,
        segments=segments,
        warnings=warnings,
        response_count=_estimate_tabular_response_count(source.asset_type, segments),
        question_scores=[],
    )


def _extract_excel_file(
    path: Path,
    source: SourceAsset,
    instructor_name: str,
) -> tuple[list[RawTextSegment], list[str]]:
    sheets = _read_excel_sheets(path)
    if not sheets:
        raise ValueError(
            f"{source.label}: 엑셀 파일이 비어 있거나 읽을 수 있는 시트를 찾지 못했습니다. "
            "응답이 담긴 단일 시트를 남기거나 CSV로 저장해 다시 업로드해 주세요."
        )

    candidate = _select_voc_sheet_candidate(sheets, source.label)
    return _build_excel_segments(
        source=source,
        instructor_name=instructor_name,
        sheet_candidate=candidate,
    )


def _extract_voc_excel_file(
    path: Path,
    source: SourceAsset,
    instructor_name: str,
) -> VocExtractionResult:
    sheets = _read_excel_sheets(path)
    if not sheets:
        raise ValueError(
            f"{source.label}: 엑셀 파일이 비어 있거나 읽을 수 있는 시트를 찾지 못했습니다. "
            "응답이 담긴 단일 시트를 남기거나 CSV로 저장해 다시 업로드해 주세요."
        )

    survey_candidate = _select_voc_survey_candidate(sheets, source.label)
    if survey_candidate is not None:
        return _build_voc_extraction_from_survey_candidate(
            source=source,
            instructor_name=instructor_name,
            survey_candidate=survey_candidate,
        )

    candidate = _select_voc_sheet_candidate(sheets, source.label)
    segments, warnings = _build_excel_segments(
        source=source,
        instructor_name=instructor_name,
        sheet_candidate=candidate,
    )
    return VocExtractionResult(
        source=source,
        segments=segments,
        warnings=warnings,
        response_count=_estimate_tabular_response_count(source.asset_type, segments),
        question_scores=[],
    )


def _read_excel_sheets(path: Path) -> list[TabularSheet]:
    suffix = path.suffix.lower()
    if suffix == ".xlsx":
        return _read_xlsx_sheets(path)
    if suffix == ".xls":
        return _read_xls_sheets(path)
    return []


def _read_xlsx_sheets(path: Path) -> list[TabularSheet]:
    from openpyxl import load_workbook

    workbook = load_workbook(filename=path, read_only=True, data_only=True)
    try:
        sheets: list[TabularSheet] = []
        for worksheet in workbook.worksheets:
            if getattr(worksheet, "sheet_state", "visible") != "visible":
                continue
            rows: list[list[str]] = []
            for row in worksheet.iter_rows(values_only=True):
                normalized_row = _normalize_sheet_row(row)
                if normalized_row:
                    rows.append(normalized_row)
            if rows:
                sheets.append(TabularSheet(name=worksheet.title, rows=rows))
        return sheets
    finally:
        workbook.close()


def _read_xls_sheets(path: Path) -> list[TabularSheet]:
    import xlrd

    workbook = xlrd.open_workbook(path)
    sheets: list[TabularSheet] = []
    for worksheet in workbook.sheets():
        rows: list[list[str]] = []
        for row_index in range(worksheet.nrows):
            normalized_row = _normalize_sheet_row(
                worksheet.cell_value(row_index, column_index)
                for column_index in range(worksheet.ncols)
            )
            if normalized_row:
                rows.append(normalized_row)
        if rows:
            sheets.append(TabularSheet(name=worksheet.name, rows=rows))
    return sheets


def _normalize_sheet_row(values) -> list[str]:
    row = [_normalize_tabular_cell(value) for value in values]
    while row and not row[-1]:
        row.pop()
    return row if any(row) else []


def _normalize_tabular_cell(value) -> str:  # noqa: ANN001
    if value is None:
        return ""
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    if isinstance(value, datetime):
        return value.isoformat(sep=" ", timespec="minutes")
    if isinstance(value, date):
        return value.isoformat()
    if isinstance(value, float):
        if value.is_integer():
            return str(int(value))
        return normalize_text(f"{value:.6f}".rstrip("0").rstrip("."))
    if isinstance(value, int):
        return str(value)
    return normalize_text(str(value))


def _select_voc_sheet_candidate(sheets: list[TabularSheet], source_label: str) -> VocSheetCandidate:
    candidates = [
        candidate
        for sheet in sheets
        if (candidate := _score_voc_sheet_candidate(sheet)) is not None
    ]
    if not candidates:
        raise ValueError(
            f"{source_label}: 엑셀 구조가 모호합니다. 응답이 담긴 단일 시트만 남기거나 CSV로 저장해 다시 업로드해 주세요."
        )

    candidates.sort(key=lambda item: item.score, reverse=True)
    top = candidates[0]
    if len(candidates) == 1:
        return top

    second = candidates[1]
    if top.score >= second.score + 1.25:
        return top

    raise ValueError(
        f"{source_label}: 엑셀 구조가 모호합니다. 응답이 담긴 단일 시트만 남기거나 CSV로 저장해 다시 업로드해 주세요."
    )


def _select_voc_survey_candidate(
    sheets: list[TabularSheet],
    source_label: str,
) -> VocSurveyCandidate | None:
    candidates = [
        candidate
        for sheet in sheets
        if (candidate := _score_voc_survey_sheet_candidate(sheet)) is not None
    ]
    if not candidates:
        return None

    candidates.sort(key=lambda item: item.score, reverse=True)
    top = candidates[0]
    if len(candidates) == 1:
        return top

    second = candidates[1]
    if top.score >= second.score + 1.25:
        return top

    raise ValueError(
        f"{source_label}: 엑셀 구조가 모호합니다. 응답이 담긴 단일 시트만 남기거나 CSV로 저장해 다시 업로드해 주세요."
    )


def _score_voc_sheet_candidate(sheet: TabularSheet) -> VocSheetCandidate | None:
    header = list(sheet.rows[0]) if sheet.rows else []
    data_rows = [row for row in sheet.rows[1:] if any(cell for cell in row)] if len(sheet.rows) > 1 else []
    if not header or len(data_rows) < 2:
        return None

    column_count = max(len(header), max((len(row) for row in data_rows), default=0))
    header = _pad_row(header, column_count)
    padded_rows = [_pad_row(row, column_count) for row in data_rows]
    non_empty_cells = 0
    numeric_like_cells = 0
    text_columns: list[int] = []
    meta_columns: list[int] = []
    header_hint_count = 0

    for column_index in range(column_count):
        values = [row[column_index] for row in padded_rows if row[column_index]]
        if not values:
            continue
        non_empty_cells += len(values)
        numeric_like_cells += sum(1 for value in values if _looks_numeric_like(value))
        header_value = normalize_text(header[column_index]).lower()
        if _contains_hint(header_value, EXCEL_VOC_HEADER_HINTS):
            header_hint_count += 1
        if _is_text_rich_column(values, header_value):
            text_columns.append(column_index)
        elif _contains_hint(header_value, EXCEL_VOC_META_HEADER_HINTS):
            meta_columns.append(column_index)

    numeric_ratio = (numeric_like_cells / non_empty_cells) if non_empty_cells else 1.0
    response_like_rows = sum(
        1
        for row in padded_rows
        if any(column_index < len(row) and _is_text_rich_value(row[column_index]) for column_index in text_columns)
    )
    if not text_columns or response_like_rows < 2 or numeric_ratio > 0.82:
        return None

    score = (
        min(len(padded_rows), 50) * 0.08
        + len(text_columns) * 1.2
        + len(meta_columns) * 0.2
        + response_like_rows * 0.15
        + header_hint_count * 0.6
        - numeric_ratio * 1.5
    )

    return VocSheetCandidate(
        sheet=sheet,
        score=score,
        header=header,
        data_rows=padded_rows,
        text_columns=text_columns,
        meta_columns=meta_columns,
        response_like_rows=response_like_rows,
        numeric_ratio=numeric_ratio,
        header_hint_count=header_hint_count,
    )


def _score_voc_survey_sheet_candidate(sheet: TabularSheet) -> VocSurveyCandidate | None:
    if len(sheet.rows) < 2:
        return None

    best_candidate: VocSurveyCandidate | None = None
    max_header_index = min(4, len(sheet.rows) - 2)
    for header_row_index in range(max_header_index + 1):
        candidate = _build_voc_survey_candidate(sheet, header_row_index)
        if candidate is None:
            continue
        if best_candidate is None or candidate.score > best_candidate.score:
            best_candidate = candidate
    return best_candidate


def _build_voc_survey_candidate(
    sheet: TabularSheet,
    header_row_index: int,
) -> VocSurveyCandidate | None:
    if header_row_index >= len(sheet.rows) - 1:
        return None
    column_count = max(len(row) for row in sheet.rows[: header_row_index + 2])
    header = _collapse_sheet_headers(sheet.rows[: header_row_index + 1], column_count)
    leaf_header = _pad_row(sheet.rows[header_row_index], column_count)
    if not any(
        _extract_voc_question_id(value) or _is_voc_free_text_header(value)
        for value in leaf_header
    ):
        return None
    data_rows = [_pad_row(row, column_count) for row in sheet.rows[header_row_index + 1 :] if any(row)]
    if not data_rows:
        return None

    question_columns: list[VocSurveyQuestionColumn] = []
    free_text_columns: list[int] = []
    for column_index in range(column_count):
        header_value = normalize_text(header[column_index])
        if not header_value:
            continue
        values = [
            row[column_index]
            for row in data_rows
            if column_index < len(row) and normalize_text(row[column_index])
        ]
        if not values:
            continue

        question_id = _extract_voc_question_id(header_value)
        if question_id and question_id.upper().startswith("BQ"):
            numeric_values = [
                parsed
                for value in values
                if (parsed := _parse_voc_question_score_value(value)) is not None
            ]
            if len(numeric_values) >= max(2, math.ceil(len(values) * 0.65)):
                question_columns.append(
                    VocSurveyQuestionColumn(
                        column_index=column_index,
                        question_id=question_id.upper(),
                        group=_question_group(question_id),
                        label=_question_label(
                            question_id=question_id,
                            leaf_header=leaf_header[column_index],
                            collapsed_header=header_value,
                        ),
                        scale_max=max(5, int(max(numeric_values))),
                    )
                )
                continue

        if _is_voc_free_text_header(header_value) and any(_is_text_rich_value(value) for value in values):
            free_text_columns.append(column_index)

    if not question_columns:
        return None

    response_row_count = sum(
        1
        for row in data_rows
        if _row_has_voc_survey_content(
            row=row,
            question_columns=question_columns,
            free_text_columns=free_text_columns,
        )
    )
    if response_row_count < 2:
        return None

    score = (
        min(response_row_count, 120) * 0.05
        + len(question_columns) * 1.4
        + len(free_text_columns) * 0.8
        + header_row_index * 0.35
    )
    return VocSurveyCandidate(
        sheet=sheet,
        score=score,
        header_row_index=header_row_index,
        header=header,
        data_rows=data_rows,
        question_columns=question_columns,
        free_text_columns=free_text_columns,
        response_row_count=response_row_count,
    )


def _build_excel_segments(
    *,
    source: SourceAsset,
    instructor_name: str,
    sheet_candidate: VocSheetCandidate,
) -> tuple[list[RawTextSegment], list[str]]:
    included_columns = list(dict.fromkeys([*sheet_candidate.meta_columns, *sheet_candidate.text_columns]))
    if not included_columns:
        raise ValueError(
            f"{source.label}: 엑셀에서 VOC 응답 열을 찾지 못했습니다. 응답이 담긴 단일 시트만 남기거나 CSV로 저장해 다시 업로드해 주세요."
        )

    sheet_name = normalize_text(sheet_candidate.sheet.name) or "sheet"
    segments: list[RawTextSegment] = []
    for row_index, row in enumerate(sheet_candidate.data_rows, start=1):
        parts: list[str] = []
        for column_index in included_columns:
            value = row[column_index] if column_index < len(row) else ""
            if not value:
                continue
            header_value = sheet_candidate.header[column_index] if column_index < len(sheet_candidate.header) else ""
            label = normalize_text(header_value) or f"column {column_index + 1}"
            parts.append(f"{label}: {value}")
        row_text = normalize_text(" | ".join(parts))
        if not row_text:
            continue
        segments.append(
            RawTextSegment(
                source_id=source.id,
                instructor_name=instructor_name,
                source_label=source.label,
                source_type=source.asset_type,
                locator=f"{sheet_name}.row.{row_index}",
                text=row_text,
            )
        )

    if not segments:
        raise ValueError(
            f"{source.label}: 엑셀에서 분석 가능한 VOC 응답을 찾지 못했습니다. 응답이 담긴 단일 시트만 남기거나 CSV로 저장해 다시 업로드해 주세요."
        )

    warnings: list[str] = []
    if len(sheet_candidate.data_rows) != len(segments):
        warnings.append(
            f"{source.label}: 선택된 시트 '{sheet_candidate.sheet.name}'에서 비어 있는 응답 행 {len(sheet_candidate.data_rows) - len(segments)}개를 제외했습니다."
        )
    return segments, warnings


def _build_voc_extraction_from_survey_candidate(
    *,
    source: SourceAsset,
    instructor_name: str,
    survey_candidate: VocSurveyCandidate,
) -> VocExtractionResult:
    sheet_name = normalize_text(survey_candidate.sheet.name) or "sheet"
    question_scores = _build_voc_question_scores(survey_candidate)
    segments: list[RawTextSegment] = []
    skipped_empty_text_rows = 0

    for row_index, row in enumerate(survey_candidate.data_rows, start=1):
        if not _row_has_voc_survey_content(
            row=row,
            question_columns=survey_candidate.question_columns,
            free_text_columns=survey_candidate.free_text_columns,
        ):
            continue
        parts: list[str] = []
        for column_index in survey_candidate.free_text_columns:
            value = row[column_index] if column_index < len(row) else ""
            if not _is_text_rich_value(value):
                continue
            header_value = survey_candidate.header[column_index] if column_index < len(survey_candidate.header) else ""
            label = normalize_text(header_value) or f"column {column_index + 1}"
            parts.append(f"{label}: {normalize_text(value)}")
        row_text = normalize_text(" | ".join(parts))
        if not row_text:
            skipped_empty_text_rows += 1
            continue
        segments.append(
            RawTextSegment(
                source_id=source.id,
                instructor_name=instructor_name,
                source_label=source.label,
                source_type=source.asset_type,
                locator=f"{sheet_name}.row.{row_index}",
                text=row_text,
            )
        )

    warnings: list[str] = []
    if survey_candidate.free_text_columns and skipped_empty_text_rows:
        warnings.append(
            f"{source.label}: 선택된 시트 '{survey_candidate.sheet.name}'에서 자유의견이 비어 있는 응답 {skipped_empty_text_rows}개를 텍스트 분석에서 제외했습니다."
        )

    return VocExtractionResult(
        source=source,
        segments=segments,
        warnings=warnings,
        response_count=survey_candidate.response_row_count,
        question_scores=question_scores,
    )


def _build_voc_question_scores(survey_candidate: VocSurveyCandidate) -> list[dict]:
    scores: list[dict] = []
    for column in survey_candidate.question_columns:
        values: list[float] = []
        for row in survey_candidate.data_rows:
            value = row[column.column_index] if column.column_index < len(row) else ""
            if (parsed := _parse_voc_question_score_value(value)) is not None:
                values.append(parsed)
        if not values:
            continue
        average = round(sum(values) / len(values), 2)
        scores.append(
            {
                "question_id": column.question_id,
                "group": column.group,
                "label": column.label,
                "average": average,
                "response_count": len(values),
                "scale_max": column.scale_max,
            }
        )
    return scores


def _pad_row(row: list[str], size: int) -> list[str]:
    if len(row) >= size:
        return list(row)
    return list(row) + [""] * (size - len(row))


def _looks_numeric_like(value: str) -> bool:
    clean = normalize_text(value)
    if not clean:
        return False
    return bool(NUMERICISH_RE.fullmatch(clean))


def _is_text_rich_column(values: list[str], header_value: str) -> bool:
    text_like_values = [value for value in values if _is_text_rich_value(value)]
    if _contains_hint(header_value, EXCEL_VOC_TEXT_HEADER_HINTS):
        return bool(text_like_values)
    if len(text_like_values) < 2:
        return False
    average_length = sum(len(value) for value in text_like_values) / len(text_like_values)
    return average_length >= 8


def _is_text_rich_value(value: str) -> bool:
    clean = normalize_text(value)
    if not clean:
        return False
    if _looks_numeric_like(clean):
        return False
    return len(TEXT_SIGNAL_RE.findall(clean)) >= 4 or len(clean) >= 8


def _contains_hint(value: str, hints: tuple[str, ...]) -> bool:
    return any(hint in value for hint in hints)


def _estimate_tabular_response_count(source_type: str, segments: list[RawTextSegment]) -> int:
    if source_type in {"csv", "xlsx", "xls"}:
        return len(segments)
    if source_type == "text":
        return sum(max(1, len([line for line in segment.text.split("|") if line.strip()])) for segment in segments)
    return len(segments)


def _collapse_sheet_headers(header_rows: list[list[str]], column_count: int) -> list[str]:
    padded_rows = [_forward_fill_header_row(_pad_row(row, column_count)) for row in header_rows]
    collapsed: list[str] = []
    for column_index in range(column_count):
        parts: list[str] = []
        for row in padded_rows:
            cell = normalize_text(row[column_index])
            if not cell:
                continue
            if cell.lower() == (parts[-1].lower() if parts else ""):
                continue
            parts.append(cell)
        collapsed.append(" / ".join(parts))
    return collapsed


def _forward_fill_header_row(row: list[str]) -> list[str]:
    filled: list[str] = []
    previous = ""
    for value in row:
        clean = normalize_text(value)
        if clean:
            previous = clean
            filled.append(clean)
            continue
        filled.append(previous if previous else "")
    return filled


def _extract_voc_question_id(header_value: str) -> str | None:
    matches = VOC_SURVEY_QUESTION_ID_RE.findall(header_value or "")
    if not matches:
        return None
    return matches[-1].upper()


def _question_group(question_id: str) -> str:
    question = str(question_id or "").upper()
    return question.split("-", maxsplit=1)[0] if "-" in question else question


def _question_label(*, question_id: str, leaf_header: str, collapsed_header: str) -> str:
    preferred = normalize_text(leaf_header) or normalize_text(collapsed_header)
    label = re.sub(
        rf"^.*?\b{re.escape(str(question_id))}\b[\s\.\-:)]*",
        "",
        preferred,
        flags=re.IGNORECASE,
    )
    normalized = normalize_text(label)
    return normalized or normalize_text(preferred) or str(question_id).upper()


def _parse_voc_question_score_value(value: str) -> float | None:
    clean = normalize_text(value)
    if not clean:
        return None
    if clean.endswith("점"):
        clean = clean[:-1].strip()
    try:
        numeric = float(clean)
    except ValueError:
        return None
    if not (1.0 <= numeric <= 5.0):
        return None
    return numeric


def _is_voc_free_text_header(header_value: str) -> bool:
    return _contains_hint(str(header_value or "").lower(), VOC_SURVEY_FREE_TEXT_HINTS)


def _row_has_voc_survey_content(
    *,
    row: list[str],
    question_columns: list[VocSurveyQuestionColumn],
    free_text_columns: list[int],
) -> bool:
    for column in question_columns:
        value = row[column.column_index] if column.column_index < len(row) else ""
        if _parse_voc_question_score_value(value) is not None:
            return True
    for column_index in free_text_columns:
        value = row[column_index] if column_index < len(row) else ""
        if normalize_text(value):
            return True
    return False


def _shape_texts(shape) -> list[str]:  # noqa: ANN001
    texts: list[str] = []

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            texts.extend(_shape_texts(child))

    if getattr(shape, "has_text_frame", False):
        text = normalize_text(shape.text or "")
        if text:
            texts.append(text)

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                text = normalize_text(cell.text or "")
                if text:
                    texts.append(text)

    return texts


def _fetch_youtube_transcript(
    video_id: str,
    settings: Settings,
    *,
    storage: ObjectStorage | None = None,
):
    cache = YoutubeCache(settings, storage=storage)
    cached_transcript = cache.get_transcript(video_id=video_id, allow_stale=False)
    if cached_transcript is not None:
        return list(cached_transcript.value or []), []

    stale_transcript = cache.get_transcript(video_id=video_id, allow_stale=True)
    proxy_requested = bool(settings.youtube_scraperapi_enabled)
    http_client = None
    proxy_active = False
    try:
        http_client = build_youtube_scraperapi_http_client(settings, session_seed=f"video:{video_id}")
        proxy_active = http_client is not None
        throttle_youtube_requests(settings)
        transcript = YouTubeTranscriptApi(http_client=http_client).fetch(
            video_id,
            languages=["ko", "en", "en-US"],
        )
    except Exception as exc:  # noqa: BLE001
        logger.warning(
            "YouTube transcript fetch failed for %s (scraperapi_enabled=%s, proxy_active=%s, error=%s: %s)",
            video_id,
            proxy_requested,
            proxy_active,
            type(exc).__name__,
            _summarize_youtube_exception_for_log(exc),
        )
        if is_youtube_request_limited_error(exc):
            mark_youtube_request_limited(settings)
        if stale_transcript is not None and is_youtube_request_limited_error(exc):
            warning = f"{_video_url(video_id)}: {YOUTUBE_STALE_TRANSCRIPT_WARNING}"
            return list(stale_transcript.value or []), [warning]
        raise

    serialized = [
        {
            "text": _transcript_text(item),
            "start": _transcript_start(item),
            "duration": _transcript_duration(item),
        }
        for item in transcript
    ]
    cache.put_transcript(video_id=video_id, value=serialized)
    return serialized, []


def _summarize_youtube_exception_for_log(exc: Exception) -> str:
    message = str(exc).strip().splitlines()[0] if str(exc).strip() else type(exc).__name__
    return message if len(message) <= 160 else f"{message[:157]}..."


def _extract_video_id(url: str) -> str | None:
    direct_match = YOUTUBE_ID_RE.search(url)
    if direct_match:
        return direct_match.group(1)

    parsed = urlparse(url)
    if parsed.hostname in {"www.youtube.com", "youtube.com"}:
        return parse_qs(parsed.query).get("v", [None])[0]
    return None


def _transcript_text(item) -> str:  # noqa: ANN001
    if isinstance(item, dict):
        return str(item.get("text", ""))
    return str(getattr(item, "text", ""))


def _transcript_start(item) -> float:  # noqa: ANN001
    if isinstance(item, dict):
        return float(item.get("start", 0.0))
    return float(getattr(item, "start", 0.0))


def _transcript_duration(item) -> float:  # noqa: ANN001
    if isinstance(item, dict):
        return float(item.get("duration", 0.0))
    return float(getattr(item, "duration", 0.0))


def _video_url(video_id: str) -> str:
    return f"https://www.youtube.com/watch?v={video_id}"


def _try_youtube_stt_fallback(
    url: str,
    video_id: str,
    source: SourceAsset,
    instructor_name: str,
    settings: Settings,
    *,
    storage: ObjectStorage | None = None,
) -> tuple[list[RawTextSegment], list[str]]:
    if not settings.youtube_stt_enabled:
        return [], [f"{url}: {YOUTUBE_STT_DISABLED_WARNING}"]
    if OpenAI is None or not settings.openai_api_key:
        return [], [f"{url}: STT fallback에 필요한 OpenAI 설정이 없습니다."]

    active_storage = storage or create_object_storage(settings)
    duration_seconds = _lookup_cached_youtube_duration(url, video_id, settings, active_storage)
    if _stt_budget_exceeded(duration_seconds, settings, active_storage):
        return [], [f"{url}: {YOUTUBE_STT_BUDGET_EXCEEDED_WARNING}"]

    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            audio_path, downloaded_duration_seconds = _download_youtube_audio_for_stt(
                url,
                settings,
                destination_dir=Path(temp_dir),
            )
            if not audio_path.exists():
                return [], [f"{url}: STT fallback용 오디오를 저장하지 못했습니다."]
            if audio_path.stat().st_size > settings.youtube_stt_max_file_bytes:
                size_mb = audio_path.stat().st_size / (1024 * 1024)
                return [], [f"{url}: {YOUTUBE_STT_FILE_TOO_LARGE_WARNING} ({size_mb:.1f}MB)"]

            transcript_text = _transcribe_audio_path(audio_path, settings)
    except Exception as exc:  # noqa: BLE001
        return [], [f"{url}: STT fallback에 실패했습니다. ({type(exc).__name__}: {exc})"]

    normalized_text = normalize_text(transcript_text)
    if not normalized_text:
        return [], [f"{url}: STT fallback 결과에서 분석 가능한 텍스트를 찾지 못했습니다."]

    effective_duration_seconds = max(duration_seconds, downloaded_duration_seconds)
    if effective_duration_seconds > 0:
        _record_stt_usage(effective_duration_seconds, settings, active_storage)

    return [
        RawTextSegment(
            source_id=source.id,
            instructor_name=instructor_name,
            source_label=source.label,
            source_type=source.asset_type,
            locator="stt",
            text=normalized_text,
        )
    ], [f"{url}: {YOUTUBE_STT_FALLBACK_WARNING}"]


def _download_youtube_audio_for_stt(
    url: str,
    _settings: Settings,
    *,
    destination_dir: Path,
) -> tuple[Path, int]:
    options = {
        "quiet": True,
        "no_warnings": True,
        "noplaylist": True,
        "format": "bestaudio[ext=m4a]/bestaudio[ext=mp4]/bestaudio[ext=webm]/bestaudio[ext=mp3]/bestaudio",
        "outtmpl": str(destination_dir / "%(id)s.%(ext)s"),
        "restrictfilenames": True,
    }

    try:
        with YoutubeDL(options) as ydl:
            info = ydl.extract_info(url, download=True)
            path = _resolve_downloaded_media_path(info, destination_dir)
            duration_seconds = int((info or {}).get("duration") or 0)
            return path, duration_seconds
    except Exception as exc:  # noqa: BLE001
        raise


def _resolve_downloaded_media_path(info: dict | None, destination_dir: Path) -> Path:
    payload = info or {}
    for item in payload.get("requested_downloads") or []:
        filepath = str(item.get("filepath") or item.get("_filename") or "").strip()
        if filepath:
            path = Path(filepath)
            if path.exists():
                return path

    file_candidates = sorted(destination_dir.glob("*"))
    if file_candidates:
        return file_candidates[0]
    raise RuntimeError("yt-dlp가 오디오 파일 경로를 반환하지 않았습니다.")


def _transcribe_audio_path(audio_path: Path, settings: Settings) -> str:
    client = OpenAI(api_key=settings.openai_api_key)
    with audio_path.open("rb") as audio_file:
        response = client.audio.transcriptions.create(
            file=audio_file,
            model=settings.youtube_stt_model,
            response_format="text",
        )
    return str(response or "")


def _lookup_cached_youtube_duration(
    url: str,
    video_id: str,
    settings: Settings,
    storage: ObjectStorage,
) -> int:
    cache = YoutubeCache(settings, storage=storage)
    for candidate_url in (url, _video_url(video_id)):
        cached = cache.get_metadata(
            url=candidate_url,
            max_videos=1,
            treat_as_playlist=False,
            allow_stale=True,
        )
        if cached is None:
            continue
        try:
            return int((cached.value or {}).get("duration") or 0)
        except (TypeError, ValueError, AttributeError):
            continue
    return 0


def _lookup_cached_youtube_title(
    url: str,
    video_id: str | None,
    settings: Settings,
    storage: ObjectStorage,
) -> str:
    cache = YoutubeCache(settings, storage=storage)
    candidate_urls = [str(url or "").strip()]
    if video_id:
        candidate_urls.append(_video_url(video_id))

    for candidate_url in candidate_urls:
        if not candidate_url:
            continue
        cached = cache.get_metadata(
            url=candidate_url,
            max_videos=1,
            treat_as_playlist=False,
            allow_stale=True,
        )
        if cached is None:
            continue
        title = normalize_text(str((cached.value or {}).get("title") or ""))
        if title:
            return title
    return ""


def _stt_budget_exceeded(duration_seconds: int, settings: Settings, storage: ObjectStorage) -> bool:
    budget_minutes = max(0, int(settings.youtube_stt_monthly_minutes_budget or 0))
    if budget_minutes <= 0 or duration_seconds <= 0:
        return False
    payload = _load_stt_usage(storage)
    used_minutes = float(payload.get("used_minutes", 0.0) or 0.0)
    requested_minutes = duration_seconds / 60.0
    return (used_minutes + requested_minutes) > budget_minutes


def _record_stt_usage(duration_seconds: int, settings: Settings, storage: ObjectStorage) -> None:
    budget_minutes = max(0, int(settings.youtube_stt_monthly_minutes_budget or 0))
    if budget_minutes <= 0 or duration_seconds <= 0:
        return
    payload = _load_stt_usage(storage)
    payload["used_minutes"] = round(float(payload.get("used_minutes", 0.0) or 0.0) + (duration_seconds / 60.0), 4)
    storage.put_json(_stt_usage_key(), payload)


def _load_stt_usage(storage: ObjectStorage) -> dict:
    try:
        payload = storage.get_json(_stt_usage_key())
    except Exception:  # noqa: BLE001
        return {"month": _stt_usage_month(), "used_minutes": 0.0}
    if payload.get("month") != _stt_usage_month():
        return {"month": _stt_usage_month(), "used_minutes": 0.0}
    return payload


def _stt_usage_key() -> str:
    return f"youtube-usage/stt/{_stt_usage_month()}.json"


def _stt_usage_month() -> str:
    return time.strftime("%Y-%m")
